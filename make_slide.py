"""
make_slide.py  â€”  CX Audit Agent executive slide generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# â”€â”€ Palette â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
BG      = RGBColor(0x08, 0x0C, 0x14)
CARD    = RGBColor(0x0F, 0x17, 0x2A)
CARD2   = RGBColor(0x11, 0x18, 0x27)
DARK    = RGBColor(0x0D, 0x12, 0x1F)
BORDER  = RGBColor(0x1E, 0x29, 0x3B)
ACCENT  = RGBColor(0x63, 0x66, 0xF1)
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
WHITE   = RGBColor(0xF1, 0xF5, 0xF9)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
DIM     = RGBColor(0x47, 0x55, 0x69)
CYAN    = RGBColor(0x06, 0xB6, 0xD4)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)

# Pre-defined dark tint backgrounds (avoids runtime color math)
DIM_ACCENT = RGBColor(0x10, 0x11, 0x30)
DIM_GREEN  = RGBColor(0x03, 0x24, 0x1A)
DIM_CYAN   = RGBColor(0x02, 0x22, 0x2A)
DIM_PURPLE = RGBColor(0x1E, 0x18, 0x38)
DIM_AMBER  = RGBColor(0x2E, 0x20, 0x04)
DIM_RED    = RGBColor(0x2E, 0x0D, 0x0D)

# â”€â”€ Slide setup â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = BG


# â”€â”€ Helpers â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def rect(x, y, w, h, fill=None, border=None, border_pt=0.75, radius=False):
    kind = 5 if radius else 1  # 5=rounded rect, 1=rect
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    if radius:
        try:
            spPr = s._element.find(qn('p:spPr'))
            if spPr is not None:
                prstGeom = spPr.find(qn('a:prstGeom'))
                if prstGeom is not None:
                    avLst = prstGeom.find(qn('a:avLst'))
                    if avLst is None:
                        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
                    else:
                        for child in list(avLst):
                            avLst.remove(child)
                    gd = etree.SubElement(avLst, qn('a:gd'))
                    gd.set('name', 'adj')
                    gd.set('fmla', 'val 20000')
        except Exception:
            pass
    return s


def txt(text, x, y, w, h, size=9, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def hline(x, y, w, color=BORDER):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# HEADER  (0" â†’ 1.02")
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
rect(0, 0, 13.33, 1.02, fill=CARD)
rect(0, 0, 0.07, 1.02, fill=ACCENT)           # left accent bar
rect(0.16, 0.13, 0.30, 0.30, fill=DIM_ACCENT,  # icon bg
     border=ACCENT, border_pt=0.8, radius=True)
txt("â—ˆ", 0.16, 0.11, 0.30, 0.33, size=13, bold=True,
    color=ACCENT, align=PP_ALIGN.CENTER)

txt("AUTONOMOUS WEB CX AUDIT AGENT",
    0.58, 0.07, 8.0, 0.38, size=17, bold=True, color=WHITE)
txt("OpenAI Responses + Playwright  â€¢  Simulates complete user journeys on Indian fintech websites  â€¢  Scores 5 CX dimensions  â€¢  Generates structured reports with friction points & ranked recommendations",
    0.58, 0.51, 8.5, 0.26, size=8, color=MUTED)

# Right stat badges
badges = [
    ("REAL BROWSER",      GREEN,  DIM_GREEN,  9.12),
    ("AI DECISION ENGINE", ACCENT2, DIM_ACCENT, 10.16),
    ("MULTI-PERSONA",     AMBER,  DIM_AMBER,  11.35),
    ("ANON + LOGGED-IN",  CYAN,   DIM_CYAN,   12.37),
]
for label, border_c, fill_c, bx in badges:
    rect(bx, 0.19, 0.86, 0.27, fill=fill_c, border=border_c, border_pt=0.8, radius=True)
    txt(label, bx + 0.03, 0.21, 0.82, 0.23, size=6.2, bold=True,
        color=border_c, align=PP_ALIGN.CENTER)

hline(0, 1.02, 13.33, color=ACCENT)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# THREE COLUMNS  (1.06" â†’ 6.12")
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
COL_Y = 1.06
COL_H = 5.06
C1_X, C1_W = 0.07,  4.12
C2_X, C2_W = 4.29,  4.55
C3_X, C3_W = 8.94,  4.32

rect(C1_X, COL_Y, C1_W, COL_H, fill=CARD,  border=BORDER, border_pt=0.5)
rect(C2_X, COL_Y, C2_W, COL_H, fill=CARD2, border=BORDER, border_pt=0.5)
rect(C3_X, COL_Y, C3_W, COL_H, fill=CARD,  border=BORDER, border_pt=0.5)

# Column header strips
HSTRIP = RGBColor(0x18, 0x18, 0x3C)
rect(C1_X, COL_Y, C1_W, 0.31, fill=HSTRIP)
rect(C2_X, COL_Y, C2_W, 0.31, fill=HSTRIP)
rect(C3_X, COL_Y, C3_W, 0.31, fill=HSTRIP)
rect(C1_X, COL_Y, 0.05, 0.31, fill=GREEN)
rect(C2_X, COL_Y, 0.05, 0.31, fill=ACCENT)
rect(C3_X, COL_Y, 0.05, 0.31, fill=AMBER)

txt("â—Ž  HOW IT WORKS â€” AGENT PIPELINE",
    C1_X+0.13, COL_Y+0.05, C1_W-0.16, 0.22, size=7.8, bold=True, color=GREEN)
txt("â—ˆ  WHAT GETS AUDITED â€” 5 CX DIMENSIONS",
    C2_X+0.13, COL_Y+0.05, C2_W-0.16, 0.22, size=7.8, bold=True, color=ACCENT2)
txt("â—‰  OUTPUTS & REPORTS GENERATED",
    C3_X+0.13, COL_Y+0.05, C3_W-0.16, 0.22, size=7.8, bold=True, color=AMBER)

# â”€â”€â”€ COL 1: PIPELINE â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
steps = [
    ("01", "PERSONA INGESTION",
     "MD profiles: goal, financial literacy, device, constraints & success criteria",
     GREEN,  DIM_GREEN),
    ("02", "PLAYWRIGHT BROWSER LAUNCH",
     "Real Chromium opens target URL  â€¢  headed or headless  â€¢  configurable slow-motion",
     CYAN,   DIM_CYAN),
    ("03", "PAGE STATE EXTRACTION",
     "DOM scraped: interactive elements, URL, visible text â†’ structured prompt context",
     PURPLE, DIM_PURPLE),
    ("04", "OpenAI DECISION ENGINE",
     "Tool-use forces structured JSON: action, target, reasoning, cx_note, fallback",
     ACCENT2, DIM_ACCENT),
    ("05", "ACTION EXECUTION + CX LOOP",
     "Click / type / scroll / navigate / hover / dismiss popup  â€¢  auto loop detection",
     AMBER,  DIM_AMBER),
    ("06", "STRUCTURED CX EVALUATION",
     "Full journey log sent to OpenAI  â€¢  5-dimension scoring with step-level citations",
     RED,    DIM_RED),
    ("07", "REPORT GENERATION",
     "TL;DR + takeaways + scores + friction + P1/P2/P3 recs  â€¢  Markdown + JSON + Dashboard",
     GREEN,  DIM_GREEN),
]

PY  = COL_Y + 0.38
SH  = 0.59   # step height
PAD = 0.10

for i, (num, title, desc, color, dim_color) in enumerate(steps):
    sy = PY + i * (SH + 0.025)
    rect(C1_X+PAD, sy+0.05, 0.27, 0.27, fill=dim_color,
         border=color, border_pt=0.8, radius=True)
    txt(num, C1_X+PAD, sy+0.04, 0.27, 0.29, size=6.5, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    txt(title, C1_X+PAD+0.33, sy+0.03, C1_W-PAD-0.42, 0.19,
        size=7.5, bold=True, color=WHITE)
    txt(desc,  C1_X+PAD+0.33, sy+0.23, C1_W-PAD-0.42, 0.34,
        size=6.6, color=MUTED)
    if i < len(steps)-1:
        txt("â†“", C1_X+PAD+0.06, sy+SH+0.01, 0.16, 0.14,
            size=7, color=DIM, align=PP_ALIGN.CENTER)

# â”€â”€â”€ COL 2: CX DIMENSIONS â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
dims = [
    ("DISCOVERABILITY & INFORMATION ARCHITECTURE",
     "Can persona find product without prior site knowledge? Measures nav depth, search quality, category logic vs. literacy level.",
     ACCENT2, DIM_ACCENT, "8.3"),
    ("CONTENT QUALITY & FINANCIAL CLARITY",
     "Are rates, fees & eligibility stated upfront? Language fit for literacy level. Are hidden charges buried deep in the flow?",
     GREEN,   DIM_GREEN,  "7.1"),
    ("TRUST & CREDIBILITY SIGNALS",
     "RBI/NBFC compliance visibility, data security cues, social proof â€” does persona feel safe sharing PAN/Aadhaar/bank details?",
     CYAN,    DIM_CYAN,   "7.8"),
    ("CONVERSION & TASK FLOW DESIGN",
     "CTA clarity & prominence, step count in application funnel, form field complexity, error handling, login wall friction.",
     AMBER,   DIM_AMBER,  "6.4"),
    ("EMOTIONAL EXPERIENCE & PERSONA FIT",
     "Confidence / confusion / frustration moments unique to this persona. Tone match, digital anxiety triggers, overwhelm signals.",
     PURPLE,  DIM_PURPLE, "6.9"),
]

DY   = COL_Y + 0.38
DH   = 0.59

for i, (name, desc, color, dim_color, score) in enumerate(dims):
    dy = DY + i * (DH + 0.04)
    rect(C2_X+0.10, dy, C2_W-0.20, DH, fill=dim_color,
         border=RGBColor(0x1E, 0x29, 0x3B), border_pt=0.6)
    rect(C2_X+0.10, dy, 0.05, DH, fill=color)
    # Score badge
    rect(C2_X+C2_W-0.46, dy+0.07, 0.33, 0.22,
         fill=BG, border=color, border_pt=0.6, radius=True)
    txt(score+"/10", C2_X+C2_W-0.46, dy+0.08, 0.33, 0.22,
        size=6.5, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(name,  C2_X+0.21, dy+0.05, C2_W-0.62, 0.21,
        size=7, bold=True, color=WHITE)
    txt(desc,  C2_X+0.21, dy+0.27, C2_W-0.31, 0.30,
        size=6.5, color=MUTED)

# "Also uncovers" row at bottom
AU_Y = DY + 5*(DH+0.04) + 0.02
rect(C2_X+0.10, AU_Y, C2_W-0.20, 0.28, fill=CARD, border=BORDER, border_pt=0.5)
txt("ALSO UNCOVERS â†’", C2_X+0.16, AU_Y+0.05, 1.05, 0.20, size=6.2, bold=True, color=DIM)
txt("Login walls mid-funnel  Â·  Hidden charges  Â·  KYC/doc overload  Â·  Unclear CTAs  Â·  Mobile rendering issues",
    C2_X+1.22, AU_Y+0.05, C2_W-1.32, 0.20, size=6.5, color=MUTED)

# â”€â”€â”€ COL 3: OUTPUTS â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
OY  = COL_Y + 0.38
OP  = 0.12

# Per-persona header
rect(C3_X+OP, OY, C3_W-OP*2, 0.21,
     fill=DIM_GREEN, border=GREEN, border_pt=0.5, radius=True)
txt("PER PERSONA â€” GENERATED FOR EVERY JOURNEY",
    C3_X+OP+0.08, OY+0.04, C3_W-OP*2-0.10, 0.16,
    size=6.5, bold=True, color=GREEN)

persona_outputs = [
    ("CX Score  0â€“10",            "Weighted across all 5 dimensions with rationale",         ACCENT2),
    ("TL;DR Verdict",             "One-sentence executive summary of entire journey",          WHITE),
    ("Key Takeaways  Ã—3â€“5",       "Crisp bullets a PM cannot afford to miss",                WHITE),
    ("Friction Points",           "HIGH / MEDIUM / LOW  â€¢  location + exact user impact",    RED),
    ("P1 / P2 / P3 Recs",         "Prioritised actions with expected business impact",        AMBER),
    ("Step-by-Step Decision Log", "Full reasoning + CX observation at every agent step",      MUTED),
    ("Screenshots per step",      "Browser captures at each click / navigate / type action",  MUTED),
]

for i, (label, detail, color) in enumerate(persona_outputs):
    iy = OY + 0.28 + i * 0.33
    rect(C3_X+OP, iy+0.02, 0.04, 0.22, fill=color)
    txt(label,  C3_X+OP+0.10, iy,      C3_W-OP-0.13, 0.19, size=7.5, bold=True, color=color)
    txt(detail, C3_X+OP+0.10, iy+0.18, C3_W-OP-0.13, 0.16, size=6.4, color=MUTED)

hline(C3_X+OP, OY+0.28+7*0.33+0.06, C3_W-OP*2, color=BORDER)

# Cross-persona header
CRP_Y = OY + 0.28 + 7*0.33 + 0.14
rect(C3_X+OP, CRP_Y, C3_W-OP*2, 0.21,
     fill=DIM_AMBER, border=AMBER, border_pt=0.5, radius=True)
txt("MASTER REPORT â€” ACROSS ALL PERSONAS",
    C3_X+OP+0.08, CRP_Y+0.04, C3_W-OP*2-0.10, 0.16,
    size=6.5, bold=True, color=AMBER)

cross = [
    ("Cross-persona friction themes & score distribution",   MUTED),
    ("All HIGH-severity issues consolidated in one view",     RED),
    ("All P1 recommendations ranked by impact",              AMBER),
    ("Best / worst persona journey comparison",              MUTED),
]
for i, (label, color) in enumerate(cross):
    iy = CRP_Y + 0.28 + i * 0.22
    rect(C3_X+OP, iy+0.02, 0.04, 0.17, fill=color)
    txt(label, C3_X+OP+0.10, iy, C3_W-OP-0.13, 0.21, size=7, color=color)

# Delivery
DEL_Y = CRP_Y + 0.28 + 4*0.22 + 0.08
hline(C3_X+OP, DEL_Y, C3_W-OP*2, color=BORDER)
txt("DELIVERY", C3_X+OP, DEL_Y+0.06, 0.70, 0.18, size=6.2, bold=True, color=DIM)
delivery = [
    ("ðŸ“„ Markdown Reports", "reports/{persona}/report.md",     MUTED),
    ("{} JSON Journey Logs", "reports/{persona}/journey_log.json", MUTED),
    ("â—±  Live Dashboard",   "Flask  â€¢  localhost:5000",        ACCENT2),
]
for i, (mode, detail, color) in enumerate(delivery):
    iy = DEL_Y + 0.28 + i * 0.24
    txt(mode,   C3_X+OP,        iy, 1.55, 0.22, size=7, bold=True, color=color)
    txt(detail, C3_X+OP+1.56,   iy, C3_W-OP-1.58, 0.22, size=6.4, color=DIM)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# FOOTER  (6.16" â†’ 7.45")
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
FY = 6.16
rect(0, FY, 13.33, 1.34, fill=CARD)
hline(0, FY, 13.33, color=ACCENT)

txt("TECH STACK", 0.18, FY+0.08, 0.95, 0.18, size=6.2, bold=True, color=DIM)

tech = [
    ("Python 3.11",         MUTED,   0.18),
    ("Playwright",          GREEN,   1.26),
    ("OpenAI GPT-4.1",   ACCENT2, 2.22),
    ("OpenAI SDK",       PURPLE,  3.46),
    ("Tool Use / JSON",     CYAN,    4.54),
    ("Flask + REST",        AMBER,   5.58),
    ("Markdown + JSON",     MUTED,   6.60),
    ("Live Dashboard",      ACCENT2, 7.74),
]
for label, color, tx in tech:
    w = len(label) * 0.074 + 0.18
    rect(tx, FY+0.30, w, 0.26, fill=DARK, border=color, border_pt=0.5, radius=True)
    txt(label, tx+0.05, FY+0.31, w-0.06, 0.24, size=6.8, color=color)

# Right capability chips
caps = [
    ("UP TO 40 STEPS",    "per persona",         GREEN),
    ("LOOP DETECTION",    "auto-terminates",     AMBER),
    ("AUTH MODES",        "logged in & anon",    CYAN),
    ("PARALLEL RUNS",     "multi-persona",       ACCENT2),
]
for i, (title, sub, color) in enumerate(caps):
    cx = 9.10 + i * 1.06
    rect(cx, FY+0.10, 1.00, 0.52,
         fill=DARK, border=color, border_pt=0.6, radius=True)
    txt(title, cx+0.06, FY+0.12, 0.88, 0.20, size=6.5, bold=True, color=color)
    txt(sub,   cx+0.06, FY+0.34, 0.88, 0.18, size=6,   color=MUTED)

txt("Powered by OpenAI  â€¢  Built for Indian BFSI & Fintech  â€¢  Evidence-based CX scoring  â€¢  Outputs: Markdown + JSON + Interactive Dashboard",
    0.18, FY+1.06, 11.0, 0.22, size=6.4, color=DIM)

# â”€â”€ Save â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
OUT = r"C:\Users\Hiya Bhandari\Desktop\Folders\CX Audit Agent\CX_Audit_Agent.pptx"
prs.save(OUT)
print(f"âœ…  Saved: {OUT}")

